"""
レイアウトエンジン
yaml:image ブロック付きMarkdown → PPTX / PDF
座標情報をもとに画像をページ正確な位置に配置して高再現変換を行う。
"""

import io
import os
import re
from dataclasses import dataclass
from collections import defaultdict


# ──────────────────────────────────────────
# Markdown パーサー
# ──────────────────────────────────────────

@dataclass
class SlideContent:
    """1スライド分のコンテンツ。"""
    page: int
    title: str
    body_lines: list[str]
    images: list[dict]


def parse_markdown_layout(markdown: str) -> list[SlideContent]:
    """
    yaml:image ブロック付きMarkdownをパースしてスライドリストを返す。

    ルール：
    - ## Page N → 新しいスライド開始
    - # Heading → 最初のスライドのタイトル
    - ```yaml:image ... ``` → 画像メタデータとして抽出（本文から除去）
    - その他のテキスト → 本文行

    Args:
        markdown: yaml:image ブロック付きMarkdown文字列

    Returns:
        SlideContent のリスト
    """
    import yaml

    # yaml:image ブロックを全て抽出
    image_blocks: list[dict] = []

    def collect_image(match: re.Match) -> str:
        try:
            data = yaml.safe_load(match.group(1))
            if isinstance(data, dict):
                image_blocks.append(data)
        except Exception:
            pass
        return ""  # ブロックを除去

    clean_md = re.sub(
        r"```yaml:image\n(.*?)```",
        collect_image,
        markdown,
        flags=re.DOTALL,
    )

    # ページごとに画像をグループ化
    imgs_by_page: dict[int, list[dict]] = defaultdict(list)
    for img in image_blocks:
        imgs_by_page[img.get("page", 1)].append(img)

    # スライドに分割（## Page N で区切る）
    slides: list[SlideContent] = []
    current_page = 1
    current_title = ""
    current_body: list[str] = []

    # ドキュメント全体タイトル（# 見出し）を先に抽出
    doc_title_match = re.search(r"^# (.+)$", clean_md, re.MULTILINE)
    doc_title = doc_title_match.group(1).strip() if doc_title_match else ""

    def flush():
        nonlocal current_title, current_body
        # 画像は消費したものを記録して重複配置を防ぐ
        page_imgs = imgs_by_page.pop(current_page, [])
        slides.append(SlideContent(
            page=current_page,
            title=current_title,
            body_lines=[l for l in current_body if l.strip()],
            images=page_imgs,
        ))
        current_title = ""
        current_body = []

    lines = clean_md.splitlines()
    first_page_seen = False

    for line in lines:
        # # タイトル行はスキップ（doc_title として先に抽出済み）
        if re.match(r"^# ", line):
            continue

        # ## Page N → 新スライド開始
        m = re.match(r"^##\s+[Pp]age\s+(\d+)", line)
        if m:
            if first_page_seen:
                flush()
            first_page_seen = True
            current_page = int(m.group(1))
            # 最初のページにドキュメントタイトルをあてる
            if current_page == 1 and doc_title and not current_title:
                current_title = doc_title
            continue

        # ## サブ見出し（Page N 以外）→ スライドタイトルとして扱う
        if line.startswith("## ") and not current_title:
            current_title = line[3:].strip()
            continue

        current_body.append(line)

    # 最後のスライドをフラッシュ
    if first_page_seen or current_title or current_body:
        flush()

    # 1スライドもなければ全文を1枚に
    if not slides:
        all_text = [l for l in clean_md.splitlines() if l.strip()]
        slides.append(SlideContent(
            page=1,
            title="変換結果",
            body_lines=all_text,
            images=imgs_by_page.get(1, []),
        ))

    return slides


# ──────────────────────────────────────────
# PPTX 生成
# ──────────────────────────────────────────

def markdown_to_pptx_with_layout(markdown: str, base_dir: str) -> bytes:
    """
    yaml:image ブロック付きMarkdownからPPTXを生成する。
    画像は正規化座標（0-1）をもとにスライドサイズに合わせて配置。

    Args:
        markdown: yaml:image ブロック付きMarkdown文字列
        base_dir: 画像ファイルを探すベースディレクトリ

    Returns:
        PPTXファイルのバイト列
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # ワイド16:9 スライドサイズ（標準）
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    blank_layout = prs.slide_layouts[6]  # Blank
    slides_data = parse_markdown_layout(markdown)

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # ── タイトルテキストボックス
        if slide_data.title:
            tb = slide.shapes.add_textbox(
                Emu(int(0.04 * slide_w)),
                Emu(int(0.04 * slide_h)),
                Emu(int(0.92 * slide_w)),
                Emu(int(0.14 * slide_h)),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            p.alignment = PP_ALIGN.LEFT

        # ── 本文テキストボックス（画像がある場合は左半分に収める）
        has_images = bool(slide_data.images)
        body_w_ratio = 0.45 if has_images else 0.92
        body_lines = slide_data.body_lines

        if body_lines:
            tb = slide.shapes.add_textbox(
                Emu(int(0.04 * slide_w)),
                Emu(int(0.20 * slide_h)),
                Emu(int(body_w_ratio * slide_w)),
                Emu(int(0.74 * slide_h)),
            )
            tf = tb.text_frame
            tf.word_wrap = True

            first = True
            for line in body_lines:
                text = _strip_markdown_syntax(line)
                if not text:
                    continue
                if first:
                    tf.paragraphs[0].text = text
                    tf.paragraphs[0].font.size = Pt(14)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0x3D, 0x3D, 0x55)
                    first = False
                else:
                    para = tf.add_paragraph()
                    para.text = text
                    para.font.size = Pt(14)
                    para.font.color.rgb = RGBColor(0x3D, 0x3D, 0x55)
                    # 箇条書き行はインデント
                    if line.lstrip().startswith(("- ", "* ", "• ")):
                        para.level = 1

        # ── 画像配置
        for img_data in slide_data.images:
            img_rel_path = img_data.get("path", "")
            img_path = os.path.join(base_dir, img_rel_path)
            if not os.path.exists(img_path):
                continue

            x_r = float(img_data.get("x", 0.5))
            y_r = float(img_data.get("y", 0.2))
            w_r = float(img_data.get("width", 0.45))
            h_r = float(img_data.get("height", 0.6))

            # 画像が本文テキストと重ならないよう右側にオフセット
            if has_images and body_lines:
                x_r = max(x_r, 0.50)
                w_r = min(w_r, 0.46)

            left = Emu(int(x_r * slide_w))
            top = Emu(int(y_r * slide_h))
            width = Emu(int(w_r * slide_w))
            height = Emu(int(h_r * slide_h))

            # スライド境界をはみ出さないようにクリップ
            left = _clamp_emu(left, Emu(0), slide_w - Emu(100000))
            top = _clamp_emu(top, Emu(0), slide_h - Emu(100000))
            width = _clamp_emu(width, Emu(100000), slide_w - left)
            height = _clamp_emu(height, Emu(100000), slide_h - top)

            try:
                slide.shapes.add_picture(img_path, left, top, width, height)
            except Exception:
                continue

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ──────────────────────────────────────────
# PDF 生成（ReportLab）
# ──────────────────────────────────────────

def markdown_to_pdf_with_layout(markdown: str, base_dir: str) -> bytes:
    """
    yaml:image ブロック付きMarkdownからPDFを生成する。
    ReportLabを使って画像を正確な位置に配置。

    Args:
        markdown: yaml:image ブロック付きMarkdown文字列
        base_dir: 画像ファイルを探すベースディレクトリ

    Returns:
        PDFファイルのバイト列
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.lib.units import pt
    except ImportError:
        raise ImportError(
            "reportlab がインストールされていません。\n"
            "pip install reportlab を実行してください。"
        )

    # 日本語フォント登録
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("HeiseiMin-W3"))
        jp_font = "HeiseiMin-W3"
    except Exception:
        jp_font = "Helvetica"

    page_size = landscape(A4)
    page_w, page_h = page_size

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=page_size)

    slides_data = parse_markdown_layout(markdown)

    for slide_data in slides_data:
        # タイトル描画
        if slide_data.title:
            c.setFont(jp_font, 22)
            c.setFillColorRGB(0.1, 0.1, 0.18)
            c.drawString(0.04 * page_w, page_h - 0.1 * page_h, slide_data.title)

        # 本文描画
        y_cursor = page_h - 0.20 * page_h
        c.setFont(jp_font, 11)
        c.setFillColorRGB(0.24, 0.24, 0.33)
        for line in slide_data.body_lines:
            text = _strip_markdown_syntax(line)
            if not text:
                y_cursor -= 6
                continue
            c.drawString(0.04 * page_w, y_cursor, text)
            y_cursor -= 16
            if y_cursor < 0.1 * page_h:
                break

        # 画像描画
        for img_data in slide_data.images:
            img_rel_path = img_data.get("path", "")
            img_path = os.path.join(base_dir, img_rel_path)
            if not os.path.exists(img_path):
                continue

            x_r = float(img_data.get("x", 0.5))
            y_r = float(img_data.get("y", 0.15))
            w_r = float(img_data.get("width", 0.45))
            h_r = float(img_data.get("height", 0.65))

            # ReportLabはy軸が下から上なので変換
            x = x_r * page_w
            y = page_h - (y_r + h_r) * page_h
            w = w_r * page_w
            h = h_r * page_h

            try:
                c.drawImage(img_path, x, y, width=w, height=h,
                            preserveAspectRatio=True, anchor="nw")
            except Exception:
                continue

        c.showPage()

    c.save()
    return buf.getvalue()


# ──────────────────────────────────────────
# ユーティリティ
# ──────────────────────────────────────────

def _strip_markdown_syntax(line: str) -> str:
    """Markdown記法を除去してプレーンテキストにする。"""
    line = re.sub(r"^#{1,6}\s+", "", line)
    line = re.sub(r"^\s*[\*\-\+]\s+", "", line)
    line = re.sub(r"^\s*\d+\.\s+", "", line)
    line = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", line)
    line = re.sub(r"`([^`]+)`", r"\1", line)
    return line.strip()


def _clamp_emu(value, min_val, max_val):
    """EMU値を範囲内にクランプする。"""
    return max(min_val, min(value, max_val))
