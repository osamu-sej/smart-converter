"""
python-pptx エンジン
Markdown → PPTX 変換。pandoc が失敗した場合のフォールバック。
見出し(##)をスライドタイトル、本文を箇条書きとして出力する。
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt


def markdown_to_pptx_bytes(markdown: str) -> bytes:
    """
    MarkdownをPPTXに変換してバイト列で返す。

    変換ルール：
    - # 見出し  → プレゼンタイトルスライド
    - ## 見出し → 新しいスライドのタイトル
    - 本文行    → 箇条書き
    - 表        → テキスト形式で出力

    Args:
        markdown: 変換元のMarkdown文字列

    Returns:
        PPTXファイルのバイト列
    """
    import io

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]   # タイトルスライド
    content_slide_layout = prs.slide_layouts[1]  # タイトル＋コンテンツ

    lines = markdown.splitlines()
    current_title = ""
    current_body: list[str] = []

    def flush_slide():
        """蓄積した内容でスライドを追加する。"""
        if not current_title and not current_body:
            return
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = current_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(current_body):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                tf.add_paragraph().text = line

    # H1 はプレゼンタイトルスライドとして使う
    h1_added = False

    for line in lines:
        # H1 → タイトルスライド（最初の1枚のみ）
        if line.startswith("# ") and not h1_added:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = line[2:].strip()
            h1_added = True
            continue

        # H2 → 新しいコンテンツスライド
        if line.startswith("## "):
            flush_slide()
            current_title = line[3:].strip()
            current_body = []
            continue

        # 空行・区切り線はスキップ
        if not line.strip() or line.strip() == "---":
            continue

        # 箇条書き・本文・表を本文リストに追加
        text = _clean_line(line)
        if text:
            current_body.append(text)

    flush_slide()

    # スライドが1枚もない場合は全文を1スライドに
    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "変換結果"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(markdown.splitlines()[:50]):
            if line.strip():
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _clean_line(line: str) -> str:
    """Markdown記法を除去してプレーンテキストにする。"""
    # 箇条書き記号を除去
    line = re.sub(r'^[\*\-\+]\s+', '', line)
    # 番号付きリストを除去
    line = re.sub(r'^\d+\.\s+', '', line)
    # 太字・斜体
    line = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', line)
    # インラインコード
    line = re.sub(r'`([^`]+)`', r'\1', line)
    # 表の区切り行をスキップ
    if re.match(r'^[\|\s\-:]+$', line):
        return ''
    # 表の行は | を除去
    if line.startswith('|'):
        line = re.sub(r'\|', '  ', line).strip()
    return line.strip()
